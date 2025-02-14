from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # Layout 0 is a title slide
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "IEEE Student Branch: Fostering Innovation and Entrepreneurship at ADA University"
subtitle.text = "Empowering Future Innovators\nPresented by Taghi Mammadov"

# Slide 2: Introduction
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Introduction"
content = slide_2.shapes.placeholders[1].text_frame
content.text = "IEEE Student Branch at ADA University aims to foster innovation, technology development, and entrepreneurship among students."

# Slide 3: Problem Statement with Diagram
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Problem Statement"
content = slide_3.shapes.placeholders[1].text_frame
content.text = "• Students face challenges in connecting academic learning with real-world industry needs."
p = content.add_paragraph()
p.text = "• Opportunities for students to become future engineers and entrepreneurs are limited."

# Add a shape (e.g., rectangle) to represent the problem visually
shape = slide_3.shapes.add_shape(3, Inches(5), Inches(2), Inches(3), Inches(1.5))  # Shape: Rectangle (Shape ID: 3)
text_box = shape.text_frame
text_box.text = "Problem Visualized"

# Slide 4: Our Solution with Description
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Our Solution"
content = slide_4.shapes.placeholders[1].text_frame
content.text = "• IEEE Student Branch bridges students with industry professionals and resources."
p = content.add_paragraph()
p.text = "• Provides technical workshops, hackathons, and guest speakers to encourage innovation."

# Add a shape (e.g., oval) to represent the solution visually
shape = slide_4.shapes.add_shape(9, Inches(5), Inches(2), Inches(3), Inches(1.5))  # Shape: Oval (Shape ID: 9)
text_box = shape.text_frame
text_box.text = "Solution Visualized"

# Slide 5: Market Opportunity with Description
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Market Opportunity"
content = slide_5.shapes.placeholders[1].text_frame
content.text = "• IEEE has a global reach, with members in over 160 countries."
p = content.add_paragraph()
p.text = "• IEEE's support for innovation through technical societies and student competitions."

# Add a shape (e.g., circle) to represent market opportunity
shape = slide_5.shapes.add_shape(12, Inches(5), Inches(2), Inches(2), Inches(2))  # Shape: Circle (Shape ID: 12)
text_box = shape.text_frame
text_box.text = "Global Reach Visualized"

# Slide 6: Business Model
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
title.text = "Business Model"
content = slide_6.shapes.placeholders[1].text_frame
content.text = "• Partnerships with tech companies and organizations."
p = content.add_paragraph()
p.text = "• Monetization through industry collaborations, workshops, and event sponsorships."

# Add a shape (e.g., hexagon) to represent business model
shape = slide_6.shapes.add_shape(6, Inches(5), Inches(2), Inches(2.5), Inches(2))  # Shape: Hexagon (Shape ID: 6)
text_box = shape.text_frame
text_box.text = "Business Model Visualized"

# Slide 7: Our Impact
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
title.text = "Our Impact"
content = slide_7.shapes.placeholders[1].text_frame
content.text = "• The IEEE Student Branch at ADA University has already impacted students with various technical events."
p = content.add_paragraph()
p.text = "• Growth opportunities through future events and collaborations."

# Add a shape (e.g., triangle) to represent impact
shape = slide_7.shapes.add_shape(5, Inches(5), Inches(2), Inches(2.5), Inches(2))  # Shape: Triangle (Shape ID: 5)
text_box = shape.text_frame
text_box.text = "Impact Visualized"

# Slide 8: Our Team with Description
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_8.shapes.title
title.text = "Our Team"
content = slide_8.shapes.placeholders[1].text_frame
content.text = "• Led by passionate students and faculty, the IEEE Student Branch is committed to supporting future engineers."

# Slide 9: Roadmap with Description
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_9.shapes.title
title.text = "Roadmap"
content = slide_9.shapes.placeholders[1].text_frame
content.text = "• Upcoming events include:"
p = content.add_paragraph()
p.text = "• Technical workshops"
p = content.add_paragraph()
p.text = "• Coding competitions"
p = content.add_paragraph()
p.text = "• Industry collaboration projects."

# Add a rectangle to represent roadmap visually
shape = slide_9.shapes.add_shape(3, Inches(5), Inches(2), Inches(3), Inches(1.5))  # Rectangle for roadmap
text_box = shape.text_frame
text_box.text = "Roadmap Visualized"

# Slide 10: Why Join Us
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_10.shapes.title
title.text = "Why Join Us"
content = slide_10.shapes.placeholders[1].text_frame
content.text = "• Be part of a community that fosters innovation and entrepreneurship."
p = content.add_paragraph()
p.text = "• Access up-to-date information and mentorship from industry experts."

# Slide 11: Conclusion & Call to Action
slide_11 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_11.shapes.title
title.text = "Conclusion"
content = slide_11.shapes.placeholders[1].text_frame
content.text = "• Join us in shaping the future of technology and innovation."
p = content.add_paragraph()
p.text = "• Questions? Let's discuss how we can work together!"

# Save the presentation
prs.save('IEEE_Startup_Presentation_No_Images.pptx')

print("Presentation created successfully!")